"""
PPT/PPTX 解析器

按幻灯片拆分，识别标题占位符、正文占位符、备注，
提取文本并标记重点（加粗/高亮）。
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from .models import ParsedDocument


def parse_pptx(file_path: str, priority: str = "normal") -> ParsedDocument:
    """解析 PPT/PPTX 文件，返回 ParsedDocument"""
    doc = ParsedDocument(
        file_path=file_path,
        file_type="pptx",
        priority=priority,
    )

    try:
        prs = Presentation(file_path)
    except Exception as e:
        doc.parse_errors.append(f"无法打开 PPTX 文件: {e}")
        return doc

    slide_count = len(prs.slides)
    doc.total_pages = slide_count

    for slide_num, slide in enumerate(prs.slides, start=1):
        # 按顺序处理 shape（标题 → 正文 → 备注）
        for shape in slide.shapes:
            if shape.is_placeholder:
                parsed = _parse_placeholder(shape, slide_num)
                for p in parsed:
                    doc.paragraphs.append(p)
            elif shape.has_text_frame:
                for run_info in _extract_text_runs(shape.text_frame, slide_num):
                    doc.paragraphs.append(run_info)

        # 处理备注
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                doc.add_paragraph(
                    notes_text, level=0, slide=slide_num,
                    raw_style="Notes"
                )

    return doc


def _parse_placeholder(shape, slide_num: int) -> list:
    """解析占位符，返回段落列表"""
    from .models import Paragraph

    results = []
    ph = shape.placeholder_format
    ph_type = str(ph.type) if ph.type is not None else "BODY"

    # 标题占位符 → level 1
    is_title = "TITLE" in ph_type.upper() or "CENTER" in ph_type.upper()
    # 副标题 → level 2
    is_subtitle = "SUBTITLE" in ph_type.upper()
    level = 1 if is_title else (2 if is_subtitle else 0)

    if shape.has_text_frame:
        for run_info in _extract_text_runs(shape.text_frame, slide_num):
            run_info.level = max(run_info.level, level)
            results.append(run_info)

    return results


def _extract_text_runs(text_frame, slide_num: int) -> list:
    """从 text_frame 提取段落，识别加粗/高亮等格式"""
    from .models import Paragraph

    results = []
    for para in text_frame.paragraphs:
        full_text = ""
        is_important = False

        for run in para.runs:
            full_text += run.text
            if run.font.bold or run.font.italic:
                is_important = True
            # 检测高亮（python-pptx 对高亮支持有限，走颜色检测）
            try:
                if run.font.color and run.font.color.rgb:
                    # 亮色系 → 可能为高亮标记
                    pass
            except Exception:
                pass

        if full_text.strip():
            results.append(Paragraph(
                text=full_text.strip(),
                level=para.level or 0,
                is_important=is_important,
                slide=slide_num,
                raw_style=str(para.level),
            ))

    return results
